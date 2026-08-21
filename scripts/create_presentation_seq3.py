"""
Script de génération de la présentation PowerPoint (PPTX) pour la Séquence 3 D-CLIC.
Format : 16:9 Widescreen (13.333" x 7.5")
Thème : Gestion d'une campagne marketing (Modules M5 & M6)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 1. Palette de Couleurs D-CLIC
NAVY = RGBColor(26, 58, 92)        # #1A3A5C (Primaire)
ORANGE = RGBColor(232, 145, 45)     # #E8912D (Accent)
DARK_SLATE = RGBColor(44, 62, 80)   # #2C3E50 (Texte principal)
LIGHT_SLATE = RGBColor(90, 106, 122)# #5A6A7A (Texte secondaire)
BG_LIGHT = RGBColor(245, 247, 250)  # #F5F7FA (Fond de carte)
BG_WHITE = RGBColor(255, 255, 255)  # Blanc
BORDER_COLOR = RGBColor(220, 227, 234) # #DCE3EA
TEAL = RGBColor(13, 148, 136)       # #0D9488 (Accent vert/teal)
BLUE_ACCENT = RGBColor(30, 136, 229)# #1E88E5

PROJECT_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
LOGO_PATH = os.path.join(PROJECT_DIR, "Logo.jpg")
OUTPUT_PATH = os.path.join(PROJECT_DIR, "Presentation_Sequence_3_DCLIC.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

def add_header(slide, title_text, category_text="D-CLIC • MARKETING NUMÉRIQUE • SÉQUENCE 3"):
    """Ajoute l'en-tête standard sur les diapositives de contenu."""
    # Bande supérieure
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = NAVY
    top_bar.line.fill.background()
    
    # Ligne d'accent orange
    orange_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.15), Inches(13.333), Inches(0.06))
    orange_line.fill.solid()
    orange_line.fill.fore_color.rgb = ORANGE
    orange_line.line.fill.background()
    
    # Catégorie (petit surtitre)
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(9.5), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ORANGE
    p_cat.font.name = "Segoe UI"
    
    # Titre principal
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.42), Inches(9.5), Inches(0.65))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = BG_WHITE
    p_title.font.name = "Segoe UI"
    
    # Logo en haut à droite
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(11.0), Inches(0.18), width=Inches(1.6))

def add_footer(slide, slide_num, total_slides=9):
    """Ajoute le pied de page discret."""
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.733), Inches(0.35))
    tf = footer_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = f"Projet D-CLIC (OIF) — Session Juillet 2026 | Tuteur : Jesse Adirigno Ogoula                       Diapositive {slide_num}/{total_slides}"
    p.font.size = Pt(9.5)
    p.font.color.rgb = LIGHT_SLATE
    p.font.name = "Segoe UI"

def create_card(slide, left, top, width, height, bg_color=BG_LIGHT, border_color=BORDER_COLOR):
    """Crée un conteneur style carte moderne."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    else:
        card.line.fill.background()
    return card

# ==============================================================================
# SLIDE 1 : COUVERTURE / BIENVENUE
# ==============================================================================
s1 = prs.slides.add_slide(blank_layout)

# Fond Navy global
bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg1.fill.solid()
bg1.fill.fore_color.rgb = NAVY
bg1.line.fill.background()

# Bande décorative latérale gauche
bande = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
bande.fill.solid()
bande.fill.fore_color.rgb = ORANGE
bande.line.fill.background()

# Logo principal
if os.path.exists(LOGO_PATH):
    s1.shapes.add_picture(LOGO_PATH, Inches(1.2), Inches(0.8), width=Inches(3.2))

# Badge "FORMATION D-CLIC • OIF"
badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(2.2), Inches(4.5), Inches(0.4))
badge.fill.solid()
badge.fill.fore_color.rgb = ORANGE
badge.line.fill.background()
tf_badge = badge.text_frame
tf_badge.vertical_anchor = MSO_ANCHOR.MIDDLE
p_b = tf_badge.paragraphs[0]
p_b.text = "🎯 FORMATION EN MARKETING NUMÉRIQUE"
p_b.alignment = PP_ALIGN.CENTER
p_b.font.size = Pt(11)
p_b.font.bold = True
p_b.font.color.rgb = BG_WHITE
p_b.font.name = "Segoe UI"

# Grand Titre
t_box1 = s1.shapes.add_textbox(Inches(1.2), Inches(2.8), Inches(11.0), Inches(2.0))
tf1 = t_box1.text_frame
tf1.word_wrap = True
p1 = tf1.paragraphs[0]
p1.text = "Séquence 3 : Gestion d'une Campagne Marketing"
p1.font.size = Pt(36)
p1.font.bold = True
p1.font.color.rgb = BG_WHITE
p1.font.name = "Segoe UI"

p1_sub = tf1.add_paragraph()
p1_sub.text = "Méthodes Agiles • Diagramme de Gantt • Pilotage des Indicateurs (KPI)"
p1_sub.font.size = Pt(18)
p1_sub.font.color.rgb = RGBColor(200, 220, 245)
p1_sub.font.name = "Segoe UI"
p1_sub.space_before = Pt(12)

# Carte d'information tuteur / session
c_info = create_card(s1, Inches(1.2), Inches(5.2), Inches(11.0), Inches(1.5), bg_color=RGBColor(36, 75, 115), border_color=ORANGE)
tf_c = c_info.text_frame
tf_c.margin_left = Inches(0.3)
tf_c.margin_top = Inches(0.2)
p_c1 = tf_c.paragraphs[0]
p_c1.text = "👨‍🏫 Tuteur référent : Jesse Adirigno OGOULA"
p_c1.font.size = Pt(14)
p_c1.font.bold = True
p_c1.font.color.rgb = BG_WHITE
p_c1.font.name = "Segoe UI"

p_c2 = tf_c.add_paragraph()
p_c2.text = "📅 Session d'Accompagnement en Direct  |  Cohorte Groupe 1  |  Session Juillet 2026"
p_c2.font.size = Pt(12)
p_c2.font.color.rgb = RGBColor(220, 235, 250)
p_c2.font.name = "Segoe UI"
p_c2.space_before = Pt(6)

p_c3 = tf_c.add_paragraph()
p_c3.text = "💡 Rappel : Micro coupé à l'arrivée • Activez vos caméras • Posez vos questions dans le chat"
p_c3.font.size = Pt(11)
p_c3.font.color.rgb = ORANGE
p_c3.font.name = "Segoe UI"
p_c3.space_before = Pt(6)


# ==============================================================================
# SLIDE 2 : ORDRE DU JOUR & OBJECTIFS
# ==============================================================================
s2 = prs.slides.add_slide(blank_layout)
add_header(s2, "Ordre du Jour & Objectifs de la Séance")
add_footer(s2, 2)

# Colonne 1 : Objectifs Visés (Carte Bleue claire)
create_card(s2, Inches(0.8), Inches(1.6), Inches(4.5), Inches(5.1), bg_color=BG_LIGHT, border_color=NAVY)
obj_box = s2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(4.1), Inches(4.7))
tf_obj = obj_box.text_frame
tf_obj.word_wrap = True
tf_obj.margin_left = tf_obj.margin_top = tf_obj.margin_right = tf_obj.margin_bottom = 0

p_oh = tf_obj.paragraphs[0]
p_oh.text = "🎯 OBJECTIFS PÉDAGOGIQUES"
p_oh.font.size = Pt(14)
p_oh.font.bold = True
p_oh.font.color.rgb = NAVY
p_oh.font.name = "Segoe UI"

items_obj = [
    ("Comprendre les méthodes de projet", "Identifier la différence entre approche traditionnelle (Cascade) et Agilité en marketing."),
    ("Maîtriser le rétroplanning", "Savoir découper un projet en tâches, jalons et concevoir un Diagramme de Gantt."),
    ("Piloter avec des KPI précis", "Calculer et interpréter le CPC, CPA, CTR, ROAS pour mesurer la rentabilité d'une campagne."),
    ("Consolider les modules M5-M6", "Lever tous les blocages techniques et méthodologiques.")
]

for title, desc in items_obj:
    p_t = tf_obj.add_paragraph()
    p_t.text = f"• {title}"
    p_t.font.size = Pt(12)
    p_t.font.bold = True
    p_t.font.color.rgb = DARK_SLATE
    p_t.space_before = Pt(12)
    p_t.font.name = "Segoe UI"
    
    p_d = tf_obj.add_paragraph()
    p_d.text = desc
    p_d.font.size = Pt(10.5)
    p_d.font.color.rgb = LIGHT_SLATE
    p_d.space_before = Pt(2)
    p_d.font.name = "Segoe UI"

# Colonne 2 : Déroulé de la Session (4 Cartes Étapes)
steps = [
    ("1", "00:00 - 00:15", "Icebreaker & Bilan de la Semaine", "Sondage interactif, point d'étape sur la promotion et retours d'expérience."),
    ("2", "00:15 - 00:35", "Exposés & Peer-Learning", "Présentations pratiques par les apprenants volontaires sur les 3 thématiques clés."),
    ("3", "00:35 - 00:55", "Atelier Live : Gantt & Calcul des KPI", "Démo pratique d'un rétroplanning et analyse chiffrée d'une campagne."),
    ("4", "00:55 - 01:15", "Session Q&A & Plan d'Action", "Foire aux questions ouverte, levée des doutes et prochaines étapes.")
]

for i, (num, timing, stitle, sdesc) in enumerate(steps):
    card_top = Inches(1.6 + i * 1.3)
    create_card(s2, Inches(5.6), card_top, Inches(6.9), Inches(1.15), bg_color=BG_WHITE, border_color=BORDER_COLOR)
    
    # Badge numéro
    num_badge = s2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.8), card_top + Inches(0.2), Inches(0.75), Inches(0.75))
    num_badge.fill.solid()
    num_badge.fill.fore_color.rgb = ORANGE
    num_badge.line.fill.background()
    tf_num = num_badge.text_frame
    tf_num.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_n = tf_num.paragraphs[0]
    p_n.text = num
    p_n.alignment = PP_ALIGN.CENTER
    p_n.font.size = Pt(16)
    p_n.font.bold = True
    p_n.font.color.rgb = BG_WHITE
    
    # Texte étape
    st_box = s2.shapes.add_textbox(Inches(6.7), card_top + Inches(0.12), Inches(5.6), Inches(0.9))
    tf_st = st_box.text_frame
    tf_st.word_wrap = True
    tf_st.margin_left = tf_st.margin_top = tf_st.margin_right = tf_st.margin_bottom = 0
    
    p_st = tf_st.paragraphs[0]
    p_st.text = f"{stitle}  ({timing})"
    p_st.font.size = Pt(12)
    p_st.font.bold = True
    p_st.font.color.rgb = NAVY
    p_st.font.name = "Segoe UI"
    
    p_sd = tf_st.add_paragraph()
    p_sd.text = sdesc
    p_sd.font.size = Pt(10)
    p_sd.font.color.rgb = LIGHT_SLATE
    p_sd.space_before = Pt(3)
    p_sd.font.name = "Segoe UI"


# ==============================================================================
# SLIDE 3 : ICEBREAKER & BILAN PROMO
# ==============================================================================
s3 = prs.slides.add_slide(blank_layout)
add_header(s3, "Icebreaker Interactif & Point de Cohorte")
add_footer(s3, 3)

# Carte Icebreaker (Gauche)
create_card(s3, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1), bg_color=BG_WHITE, border_color=ORANGE)
box_ib = s3.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.0), Inches(4.5))
tf_ib = box_ib.text_frame
tf_ib.word_wrap = True
tf_ib.margin_left = tf_ib.margin_top = tf_ib.margin_right = tf_ib.margin_bottom = 0

p_ibh = tf_ib.paragraphs[0]
p_ibh.text = "⚡ ICEBREAKER (Sondage Slido / Chat)"
p_ibh.font.size = Pt(14)
p_ibh.font.bold = True
p_ibh.font.color.rgb = ORANGE
p_ibh.font.name = "Segoe UI"

p_q1 = tf_ib.add_paragraph()
p_q1.text = "❓ Question 1 : Quel est votre plus grand défi quand vous gérez un projet ?"
p_q1.font.size = Pt(11.5)
p_q1.font.bold = True
p_q1.font.color.rgb = NAVY
p_q1.space_before = Pt(14)
p_q1.font.name = "Segoe UI"

options = [
    "A) Le manque de temps et d'organisation",
    "B) Le respect des délais (deadlines)",
    "C) La coordination de l'équipe",
    "D) La mesure concrète des résultats"
]
for opt in options:
    po = tf_ib.add_paragraph()
    po.text = f"   {opt}"
    po.font.size = Pt(10.5)
    po.font.color.rgb = DARK_SLATE
    po.space_before = Pt(3)

p_q2 = tf_ib.add_paragraph()
p_q2.text = "❓ Question 2 (À main levée / Chat) :"
p_q2.font.size = Pt(11.5)
p_q2.font.bold = True
p_q2.font.color.rgb = NAVY
p_q2.space_before = Pt(14)

p_q2d = tf_ib.add_paragraph()
p_q2d.text = "« Avez-vous déjà utilisé un diagramme de Gantt ou un outil comme Trello / Asana ? »"
p_q2d.font.size = Pt(10.5)
p_q2d.font.color.rgb = DARK_SLATE
p_q2d.space_before = Pt(3)

# Carte Bilan & Conseils (Droite)
create_card(s3, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.1), bg_color=BG_LIGHT, border_color=NAVY)
box_b = s3.shapes.add_textbox(Inches(7.1), Inches(1.9), Inches(5.1), Inches(4.5))
tf_b = box_b.text_frame
tf_b.word_wrap = True
tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0

p_bh = tf_b.paragraphs[0]
p_bh.text = "📊 POINT SUR LE RYTHME & LA MOTIVATION"
p_bh.font.size = Pt(14)
p_bh.font.bold = True
p_bh.font.color.rgb = NAVY
p_bh.font.name = "Segoe UI"

bilan_tips = [
    ("Ne restez pas bloqués", "La Séquence 3 est souvent perçue comme la plus dense. Il est normal de trouver les notions de gestion de projet plus abstraites au départ."),
    ("La règle des 15-20 minutes / jour", "Ne cherchez pas à tout faire en une seule soirée le week-end. Une régularité quotidienne est le secret pour ne jamais décrocher."),
    ("L'esprit d'équipe D-CLIC", "Profitez du groupe WhatsApp et des sessions directes pour poser TOUTES vos questions. Il n'y a aucune question bête !")
]

for b_title, b_desc in bilan_tips:
    p_bt = tf_b.add_paragraph()
    p_bt.text = f"💡 {b_title}"
    p_bt.font.size = Pt(11.5)
    p_bt.font.bold = True
    p_bt.font.color.rgb = DARK_SLATE
    p_bt.space_before = Pt(12)
    p_bt.font.name = "Segoe UI"
    
    p_bd = tf_b.add_paragraph()
    p_bd.text = b_desc
    p_bd.font.size = Pt(10)
    p_bd.font.color.rgb = LIGHT_SLATE
    p_bd.space_before = Pt(2)
    p_bd.font.name = "Segoe UI"


# ==============================================================================
# SLIDE 4 : LES MÉTHODOLOGIES DE GESTION DE PROJET
# ==============================================================================
s4 = prs.slides.add_slide(blank_layout)
add_header(s4, "Méthodes de Gestion de Projet en Marketing Digital")
add_footer(s4, 4)

# 3 Colonnes : Cascade (Waterfall), Agile / Scrum, Kanban
cols_data = [
    ("Approche Traditionnelle (Cascade)", "Séquentielle & Linéaire", [
        "Planification rigide de A à Z avant de lancer.",
        "Chaque étape doit être terminée avant de passer à la suivante (Conception ➔ Production ➔ Diffusion).",
        "⚠️ Limite : Peu flexible en cas d'imprévu ou de changement d'algorithme."
    ], NAVY),
    ("Méthodologie Agile (Scrum)", "Itérative & Flexible", [
        "Travail par cycles courts (Sprints de 1 à 2 semaines).",
        "Ajustement continu basé sur les données réelles et les retours d'audience.",
        "✅ Idéal pour le marketing digital, les campagnes Ads et les tests A/B."
    ], ORANGE),
    ("Méthode Kanban", "Visuelle & en Flux Continu", [
        "Visualisation des tâches en colonnes : À faire ➔ En cours ➔ Terminé.",
        "Limitation du travail en cours pour éviter la surcharge.",
        "✅ Parfait pour la production de contenus et le community management."
    ], BLUE_ACCENT)
]

for i, (m_title, m_sub, m_points, col_color) in enumerate(cols_data):
    left = Inches(0.8 + i * 4.0)
    create_card(s4, left, Inches(1.6), Inches(3.7), Inches(5.1), bg_color=BG_WHITE, border_color=col_color)
    
    # En-tête de carte
    head_box = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.6), Inches(3.7), Inches(0.9))
    head_box.fill.solid()
    head_box.fill.fore_color.rgb = col_color
    head_box.line.fill.background()
    tf_h = head_box.text_frame
    tf_h.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_ht = tf_h.paragraphs[0]
    p_ht.text = m_title
    p_ht.font.size = Pt(11)
    p_ht.font.bold = True
    p_ht.font.color.rgb = BG_WHITE
    p_ht.alignment = PP_ALIGN.CENTER
    
    p_hs = tf_h.add_paragraph()
    p_hs.text = m_sub
    p_hs.font.size = Pt(9.5)
    p_hs.font.color.rgb = RGBColor(240, 240, 240)
    p_hs.alignment = PP_ALIGN.CENTER
    
    # Corps de texte
    cbox = s4.shapes.add_textbox(left + Inches(0.2), Inches(2.65), Inches(3.3), Inches(3.9))
    tf_c = cbox.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
    
    for idx, pt in enumerate(m_points):
        p_pt = tf_c.paragraphs[0] if idx == 0 else tf_c.add_paragraph()
        p_pt.text = f"• {pt}"
        p_pt.font.size = Pt(10.5)
        p_pt.font.color.rgb = DARK_SLATE
        p_pt.font.name = "Segoe UI"
        if idx > 0:
            p_pt.space_before = Pt(12)


# ==============================================================================
# SLIDE 5 : LES 4 PHASES CLÉS D'UNE CAMPAGNE MARKETING
# ==============================================================================
s5 = prs.slides.add_slide(blank_layout)
add_header(s5, "Le Cycle de Vie d'une Campagne Marketing Réussie")
add_footer(s5, 5)

phases = [
    ("Phase 1", "Cadrage & Stratégie", [
        "Définition des objectifs SMART (Spécifique, Mesurable, Atteignable, Réaliste, Temporel).",
        "Identification du Persona cible (besoins, freins, canaux privilégiés).",
        "Fixation du budget global et de la proposition de valeur."
    ], NAVY),
    ("Phase 2", "Conception & Création", [
        "Élaboration du calendrier éditorial et du rétroplanning (Gantt).",
        "Création des assets : visuels Canva, copies publicitaires, landing pages.",
        "Configuration du tracking (Pixel Meta, Google Analytics UTM)."
    ], BLUE_ACCENT),
    ("Phase 3", "Lancement & Diffusion", [
        "Mise en ligne des campagnes sur les canaux choisis (SEO, Meta Ads, Google Ads, Emailing).",
        "Surveillance en direct des premiers retours et modération des commentaires.",
        "A/B Testing des annonces et visuels les plus performants."
    ], ORANGE),
    ("Phase 4", "Analyse & Optimisation", [
        "Suivi des KPI clés : Coût d'acquisition, Taux de conversion, ROAS.",
        "Réallocation du budget vers les canaux les plus rentables.",
        "Rédaction du rapport de bilan et capitalisation des apprentissages."
    ], TEAL)
]

for i, (ph_num, ph_title, ph_items, ph_col) in enumerate(phases):
    left = Inches(0.8 + i * 3.0)
    create_card(s5, left, Inches(1.6), Inches(2.8), Inches(5.1), bg_color=BG_LIGHT, border_color=BORDER_COLOR)
    
    # Badge d'étape
    pbadge = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.2), Inches(1.8), Inches(2.4), Inches(0.4))
    pbadge.fill.solid()
    pbadge.fill.fore_color.rgb = ph_col
    pbadge.line.fill.background()
    tf_pb = pbadge.text_frame
    tf_pb.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_btext = tf_pb.paragraphs[0]
    p_btext.text = f"{ph_num} : {ph_title}"
    p_btext.font.size = Pt(10)
    p_btext.font.bold = True
    p_btext.font.color.rgb = BG_WHITE
    p_btext.alignment = PP_ALIGN.CENTER
    
    # Texte des items
    p_tbox = s5.shapes.add_textbox(left + Inches(0.15), Inches(2.35), Inches(2.5), Inches(4.2))
    tf_pt = p_tbox.text_frame
    tf_pt.word_wrap = True
    tf_pt.margin_left = tf_pt.margin_top = tf_pt.margin_right = tf_pt.margin_bottom = 0
    
    for idx, item in enumerate(ph_items):
        p_it = tf_pt.paragraphs[0] if idx == 0 else tf_pt.add_paragraph()
        p_it.text = f"✔ {item}"
        p_it.font.size = Pt(9.5)
        p_it.font.color.rgb = DARK_SLATE
        p_it.font.name = "Segoe UI"
        if idx > 0:
            p_it.space_before = Pt(10)


# ==============================================================================
# SLIDE 6 : ATELIER LIVE : LE DIAGRAMME DE GANTT
# ==============================================================================
s6 = prs.slides.add_slide(blank_layout)
add_header(s6, "Atelier Pratique : Piloter avec le Diagramme de Gantt")
add_footer(s6, 6)

# Carte Théorie Gantt (Gauche)
create_card(s6, Inches(0.8), Inches(1.6), Inches(4.8), Inches(5.1), bg_color=BG_WHITE, border_color=NAVY)
g_box = s6.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(4.4), Inches(4.6))
tf_g = g_box.text_frame
tf_g.word_wrap = True
tf_g.margin_left = tf_g.margin_top = tf_g.margin_right = tf_g.margin_bottom = 0

p_gh = tf_g.paragraphs[0]
p_gh.text = "📌 LES 4 COMPOSANTS CLÉS DU GANTT"
p_gh.font.size = Pt(13)
p_gh.font.bold = True
p_gh.font.color.rgb = NAVY
p_gh.font.name = "Segoe UI"

g_points = [
    ("Les Tâches & Activités", "Actions précises et concrètes à mener (ex: rédiger les 5 posts, configurer la campagne Facebook)."),
    ("Les Jalons (Milestones)", "Dates clés / points d'étape sans durée marquant la fin d'une phase (ex: validation du budget, lancement officiel)."),
    ("Les Dépendances", "Tâches qui ne peuvent démarrer qu'une fois la précédente achevée (ex: impossible de lancer les Ads sans les visuels validés)."),
    ("Les Responsables", "Attribution claire : qui fait quoi pour éviter les retards.")
]

for g_title, g_desc in g_points:
    p_gt = tf_g.add_paragraph()
    p_gt.text = f"• {g_title}"
    p_gt.font.size = Pt(11)
    p_gt.font.bold = True
    p_gt.font.color.rgb = DARK_SLATE
    p_gt.space_before = Pt(8)
    
    p_gd = tf_g.add_paragraph()
    p_gd.text = g_desc
    p_gd.font.size = Pt(9.5)
    p_gd.font.color.rgb = LIGHT_SLATE
    p_gd.space_before = Pt(2)

# Carte Exemple Concret de Planning (Droite)
create_card(s6, Inches(5.9), Inches(1.6), Inches(6.6), Inches(5.1), bg_color=BG_LIGHT, border_color=ORANGE)
ex_box = s6.shapes.add_textbox(Inches(6.2), Inches(1.85), Inches(6.0), Inches(4.6))
tf_ex = ex_box.text_frame
tf_ex.word_wrap = True
tf_ex.margin_left = tf_ex.margin_top = tf_ex.margin_right = tf_ex.margin_bottom = 0

p_exh = tf_ex.paragraphs[0]
p_exh.text = "🛠️ EXEMPLE : RÉTROPLANNING SUR 4 SEMAINES"
p_exh.font.size = Pt(13)
p_exh.font.bold = True
p_exh.font.color.rgb = ORANGE
p_exh.font.name = "Segoe UI"

table_tasks = [
    ("Semaine 1", "Brief client, personas, objectifs SMART et répartition budgétaire."),
    ("Semaine 2", "Rédaction des copies, création des bannières Canva & paramétrage Ads."),
    ("Semaine 3", "🚀 JALON : Lancement officiel de la campagne & A/B testing quotidien."),
    ("Semaine 4", "Clôture de la diffusion, extraction des métriques & rapport d'analyse.")
]

for sem, desc in table_tasks:
    p_st = tf_ex.add_paragraph()
    p_st.text = f"🗓️ {sem} :"
    p_st.font.size = Pt(11)
    p_st.font.bold = True
    p_st.font.color.rgb = NAVY
    p_st.space_before = Pt(8)
    
    p_sd = tf_ex.add_paragraph()
    p_sd.text = desc
    p_sd.font.size = Pt(10)
    p_sd.font.color.rgb = DARK_SLATE
    p_sd.space_before = Pt(2)

p_tools = tf_ex.add_paragraph()
p_tools.text = "💡 Outils recommandés : Google Sheets, Trello, Asana, Monday, ClickUp, Notion."
p_tools.font.size = Pt(10)
p_tools.font.bold = True
p_tools.font.color.rgb = TEAL
p_tools.space_before = Pt(14)


# ==============================================================================
# SLIDE 7 : LES INDICATEURS CLÉS DE PERFORMANCE (KPI)
# ==============================================================================
s7 = prs.slides.add_slide(blank_layout)
add_header(s7, "Mesurer l'Efficacité : Les KPI Clés d'une Campagne")
add_footer(s7, 7)

kpi_categories = [
    ("1. Notoriété & Visibilité", [
        ("Portée (Reach)", "Nombre unique de personnes ayant vu la publication."),
        ("Impressions", "Nombre total d'affichages du message."),
        ("Coût pour Mille (CPM)", "Coût pour 1 000 affichages publicitaires.")
    ], NAVY),
    ("2. Engagement & Trafic", [
        ("Taux de Clic (CTR)", "(Clics / Impressions) × 100 — Mesure l'attractivité du message."),
        ("Coût par Clic (CPC)", "Budget dépensé / Nombre de clics obtenus."),
        ("Taux d'Engagement", "(Interactions / Portée) × 100.")
    ], ORANGE),
    ("3. Conversion & Rentabilité", [
        ("Taux de Conversion", "(Ventes ou Leads / Clics) × 100."),
        ("Coût par Acquisition (CPA)", "Budget dépensé / Nombre de clients acquis."),
        ("ROAS (Return On Ad Spend)", "(Chiffre d'affaires généré / Budget Ads) × 100.")
    ], TEAL)
]

for i, (cat_title, kpi_list, cat_color) in enumerate(kpi_categories):
    left = Inches(0.8 + i * 4.0)
    create_card(s7, left, Inches(1.6), Inches(3.7), Inches(5.1), bg_color=BG_WHITE, border_color=cat_color)
    
    # En-tête de catégorie
    head_kpi = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.6), Inches(3.7), Inches(0.65))
    head_kpi.fill.solid()
    head_kpi.fill.fore_color.rgb = cat_color
    head_kpi.line.fill.background()
    tf_hk = head_kpi.text_frame
    tf_hk.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_hkt = tf_hk.paragraphs[0]
    p_hkt.text = cat_title
    p_hkt.font.size = Pt(11.5)
    p_hkt.font.bold = True
    p_hkt.font.color.rgb = BG_WHITE
    p_hkt.alignment = PP_ALIGN.CENTER
    
    # Contenu des KPI
    k_box = s7.shapes.add_textbox(left + Inches(0.2), Inches(2.35), Inches(3.3), Inches(4.2))
    tf_k = k_box.text_frame
    tf_k.word_wrap = True
    tf_k.margin_left = tf_k.margin_top = tf_k.margin_right = tf_k.margin_bottom = 0
    
    for idx, (k_name, k_desc) in enumerate(kpi_list):
        p_kn = tf_k.paragraphs[0] if idx == 0 else tf_k.add_paragraph()
        p_kn.text = f"📈 {k_name}"
        p_kn.font.size = Pt(11)
        p_kn.font.bold = True
        p_kn.font.color.rgb = DARK_SLATE
        if idx > 0:
            p_kn.space_before = Pt(10)
        
        p_kd = tf_k.add_paragraph()
        p_kd.text = k_desc
        p_kd.font.size = Pt(9.5)
        p_kd.font.color.rgb = LIGHT_SLATE
        p_kd.space_before = Pt(2)


# ==============================================================================
# SLIDE 8 : QUESTIONS / RÉPONSES & DÉBLOCAGE
# ==============================================================================
s8 = prs.slides.add_slide(blank_layout)
add_header(s8, "Session Questions / Réponses & Entraide")
add_footer(s8, 8)

# Carte Q&A Interactive (Centre)
create_card(s8, Inches(1.5), Inches(1.6), Inches(10.333), Inches(5.1), bg_color=BG_LIGHT, border_color=BLUE_ACCENT)

qa_box = s8.shapes.add_textbox(Inches(1.9), Inches(1.9), Inches(9.5), Inches(4.5))
tf_qa = qa_box.text_frame
tf_qa.word_wrap = True
tf_qa.margin_left = tf_qa.margin_top = tf_qa.margin_right = tf_qa.margin_bottom = 0

p_qah = tf_qa.paragraphs[0]
p_qah.text = "💬 LA PAROLE EST À VOUS !"
p_qah.font.size = Pt(18)
p_qah.font.bold = True
p_qah.font.color.rgb = NAVY
p_qah.font.name = "Segoe UI"

qa_items = [
    ("Questions sur les cours M5 & M6", "Des incompréhensions sur les modules, les quiz ou les méthodes de calcul ?"),
    ("Difficultés d'accès ou d'organisation", "Besoin d'aide pour rattraper le retard ou aménager votre temps de travail quotidien ?"),
    ("Partage d'expériences pratiques", "Comment appliquez-vous ces notions dans votre activité professionnelle ou vos projets personnels ?")
]

for q_title, q_desc in qa_items:
    p_qt = tf_qa.add_paragraph()
    p_qt.text = f"👉 {q_title}"
    p_qt.font.size = Pt(13)
    p_qt.font.bold = True
    p_qt.font.color.rgb = DARK_SLATE
    p_qt.space_before = Pt(14)
    p_qt.font.name = "Segoe UI"
    
    p_qd = tf_qa.add_paragraph()
    p_qd.text = f"     {q_desc}"
    p_qd.font.size = Pt(11)
    p_qd.font.color.rgb = LIGHT_SLATE
    p_qd.space_before = Pt(3)
    p_qd.font.name = "Segoe UI"

p_cons = tf_qa.add_paragraph()
p_cons.text = "🎤 Levez la main pour prendre la parole au micro ou écrivez directement dans le chat !"
p_cons.font.size = Pt(12)
p_cons.font.bold = True
p_cons.font.color.rgb = ORANGE
p_cons.space_before = Pt(18)


# ==============================================================================
# SLIDE 9 : PLAN D'ACTION & PROCHAINES ÉTAPES
# ==============================================================================
s9 = prs.slides.add_slide(blank_layout)
add_header(s9, "Plan d'Action & Prochaines Échéances")
add_footer(s9, 9)

# Colonne 1 : Checklist de la semaine
create_card(s9, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1), bg_color=BG_WHITE, border_color=TEAL)
box_chk = s9.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.0), Inches(4.5))
tf_chk = box_chk.text_frame
tf_chk.word_wrap = True
tf_chk.margin_left = tf_chk.margin_top = tf_chk.margin_right = tf_chk.margin_bottom = 0

p_chkh = tf_chk.paragraphs[0]
p_chkh.text = "✅ VOTRE TO-DO LIST DE LA SEMAINE"
p_chkh.font.size = Pt(14)
p_chkh.font.bold = True
p_chkh.font.color.rgb = TEAL
p_chkh.font.name = "Segoe UI"

tasks_list = [
    ("Finaliser la Séquence 3", "Compléter et valider les modules M5 et M6 sur la plateforme."),
    ("Pratiquer sur un outil", "Créer un premier tableau Kanban sur Trello ou un rétroplanning sur Google Sheets."),
    ("Rattraper les éventuels retards", "Prendre 15 min/jour pour valider les activités restantes des Séquences 1 & 2."),
    ("Rester actif sur WhatsApp", "Partager vos questions et aider vos pairs.")
]

for t_title, t_desc in tasks_list:
    p_tt = tf_chk.add_paragraph()
    p_tt.text = f"☑ {t_title}"
    p_tt.font.size = Pt(11.5)
    p_tt.font.bold = True
    p_tt.font.color.rgb = DARK_SLATE
    p_tt.space_before = Pt(10)
    
    p_td = tf_chk.add_paragraph()
    p_td.text = t_desc
    p_td.font.size = Pt(10)
    p_td.font.color.rgb = LIGHT_SLATE
    p_td.space_before = Pt(2)

# Colonne 2 : Canaux & Message de fin
create_card(s9, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.1), bg_color=BG_LIGHT, border_color=NAVY)
box_end = s9.shapes.add_textbox(Inches(7.1), Inches(1.9), Inches(5.1), Inches(4.5))
tf_end = box_end.text_frame
tf_end.word_wrap = True
tf_end.margin_left = tf_end.margin_top = tf_end.margin_right = tf_end.margin_bottom = 0

p_endh = tf_end.paragraphs[0]
p_endh.text = "🤝 VOTRE TUTEUR RESTE À VOTRE ÉCOUTE"
p_endh.font.size = Pt(14)
p_endh.font.bold = True
p_endh.font.color.rgb = NAVY
p_endh.font.name = "Segoe UI"

contacts = [
    ("Groupe WhatsApp Promotion", "Pour les échanges collectifs quotidiens et les partages de ressources."),
    ("Message Privé / Mail", "Pour les difficultés individuelles, blocages ou réajustement de rythme."),
    ("Prochaine Visio d'Accompagnement", "Consultez le planning et l'invitation Google Calendar transmise par mail.")
]

for c_title, c_desc in contacts:
    p_ct = tf_end.add_paragraph()
    p_ct.text = f"📱 {c_title}"
    p_ct.font.size = Pt(11.5)
    p_ct.font.bold = True
    p_ct.font.color.rgb = DARK_SLATE
    p_ct.space_before = Pt(10)
    
    p_cd = tf_end.add_paragraph()
    p_cd.text = c_desc
    p_cd.font.size = Pt(10)
    p_cd.font.color.rgb = LIGHT_SLATE
    p_cd.space_before = Pt(2)

p_mot = tf_end.add_paragraph()
p_mot.text = "✨ « Le succès n'est que la somme de petits efforts répétés jour après jour. » Bon travail à tous !"
p_mot.font.size = Pt(10.5)
p_mot.font.bold = True
p_mot.font.color.rgb = ORANGE
p_mot.space_before = Pt(14)

# Sauvegarde
prs.save(OUTPUT_PATH)
print(f"Presentation saved successfully at: {OUTPUT_PATH}")
