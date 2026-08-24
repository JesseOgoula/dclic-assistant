from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Slide 1: Titre
slide_layout = prs.slide_layouts[0] # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Séquence 3 : Gestion d'une campagne marketing"
subtitle.text = "De la théorie à la pratique : Notions clés et exemples"

# ==========================================
# PARTIE 1 : C5 - BUDGET ET PLANNING
# ==========================================
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Module 5A : Méthodes de Gestion de Projets"
body = slide.shapes.placeholders[1].text_frame
body.text = "Définition : C'est l'application de méthodes structurées (Agile, Cycle en V) pour piloter une campagne de bout en bout."
p = body.add_paragraph()
p.text = "Pourquoi c'est important :"
p.level = 1
p = body.add_paragraph()
p.text = "Permet de respecter les délais, maîtriser les risques et coordonner les équipes de manière fluide."
p.level = 2
p = body.add_paragraph()
p.text = "Exemple :"
p.level = 1
p = body.add_paragraph()
p.text = "Utiliser la méthode 'Agile' pour lancer une campagne pub : on teste une petite annonce, on analyse, on adapte, plutôt que de tout planifier à l'avance sans flexibilité."
p.level = 2

slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Module 5B : Le Budget Marketing"
body = slide.shapes.placeholders[1].text_frame
body.text = "Définition : L'enveloppe financière allouée pour atteindre les objectifs. Il inclut les coûts publicitaires (Média), la création, et les outils."
p = body.add_paragraph()
p.text = "Notions Clés :"
p.level = 1
p = body.add_paragraph()
p.text = "Coût d'Acquisition Client (CAC), Retour sur Investissement (ROI)."
p.level = 2
p = body.add_paragraph()
p.text = "Exemple :"
p.level = 1
p = body.add_paragraph()
p.text = "Pour un budget de 5 000 € : 3 000 € pour la publicité Instagram, 1 500 € pour les créations graphiques (graphiste freelance), et 500 € pour un outil d'emailing."
p.level = 2


slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Module 5C : Le Planning (Rétroplanning)"
body = slide.shapes.placeholders[1].text_frame
body.text = "Définition : Un calendrier inversé qui part de la date de lancement prévue pour déduire les dates limites de chaque tâche (Diagramme de Gantt)."
p = body.add_paragraph()
p.text = "Pourquoi c'est important :"
p.level = 1
p = body.add_paragraph()
p.text = "Évite les retards de dernière minute et s'assure que les dépendances (ex: on ne peut pas lancer la pub si la vidéo n'est pas montée) sont gérées."
p.level = 2
p = body.add_paragraph()
p.text = "Exemple :"
p.level = 1
p = body.add_paragraph()
p.text = "Lancement le 1er Décembre (J-0). J-15 : Finaliser les visuels. J-30 : Valider le concept. J-45 : Définir le budget."
p.level = 2

# ==========================================
# PARTIE 2 : C6 - TRAVAIL EN EQUIPE
# ==========================================
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Module 6A & 6B : L'Équipe et ses Métiers"
body = slide.shapes.placeholders[1].text_frame
body.text = "Définition : Une campagne réussie repose sur une collaboration entre différents experts travaillant ensemble (organisation transversale)."
p = body.add_paragraph()
p.text = "Les Métiers Clés :"
p.level = 1
p = body.add_paragraph()
p.text = "Chef de Projet, Social Media Manager, Copywriter (Rédacteur), Graphiste, Data Analyst, Traffic Manager."
p.level = 2
p = body.add_paragraph()
p.text = "Exemple :"
p.level = 1
p = body.add_paragraph()
p.text = "Le Copywriter écrit le texte de la publicité, le Graphiste crée l'image, et le Traffic Manager la diffuse sur Facebook Ads."
p.level = 2

slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Module 6C : La Répartition des Tâches"
body = slide.shapes.placeholders[1].text_frame
body.text = "Définition : L'attribution claire des responsabilités pour éviter les doublons ou les oublis (souvent via une Matrice RACI : Qui Fait, Qui Valide, Qui est Informé)."
p = body.add_paragraph()
p.text = "Notion Clé :"
p.level = 1
p = body.add_paragraph()
p.text = "La Matrice RACI (Responsible, Accountable, Consulted, Informed)."
p.level = 2
p = body.add_paragraph()
p.text = "Exemple de répartition :"
p.level = 1
p = body.add_paragraph()
p.text = "Tâche 'Validation de l'email' : Le Community Manager rédige (Fait), le Directeur Marketing valide (Valide), le service client est prévenu du lancement (Informé)."
p.level = 2


# ==========================================
# PARTIE 3 : C7 - PLAN D'ACTION ET INDICATEURS
# ==========================================
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Module 7A : Le Plan d'Action"
body = slide.shapes.placeholders[1].text_frame
body.text = "Définition : C'est la feuille de route stratégique détaillant QUI fait QUOI, COMMENT et OÙ (canaux de diffusion)."
p = body.add_paragraph()
p.text = "Pourquoi c'est important :"
p.level = 1
p = body.add_paragraph()
p.text = "Passe de la stratégie à l'opérationnel. Il structure l'exécution de la campagne."
p.level = 2
p = body.add_paragraph()
p.text = "Exemple :"
p.level = 1
p = body.add_paragraph()
p.text = "Notre plan d'action : 1. Campagne d'acquisition sur LinkedIn (B2B) + 2. Relance par Email (Newsletter) + 3. Article de blog SEO."
p.level = 2

slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Module 7B : Les Indicateurs de Suivi (KPIs)"
body = slide.shapes.placeholders[1].text_frame
body.text = "Définition : Les KPI (Key Performance Indicators) sont des métriques chiffrées qui permettent de mesurer le succès d'une campagne par rapport aux objectifs."
p = body.add_paragraph()
p.text = "Notions Clés :"
p.level = 1
p = body.add_paragraph()
p.text = "Taux de clic (CTR), Taux de conversion, Coût par Clic (CPC), Chiffre d'Affaires généré."
p.level = 2
p = body.add_paragraph()
p.text = "Exemple :"
p.level = 1
p = body.add_paragraph()
p.text = "L'objectif n'est pas 'Avoir beaucoup de j'aime', mais 'Générer 50 prospects qualifiés à un coût maximum de 10€ par prospect'."
p.level = 2

slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Module 7C : Le Tableau de Bord (Dashboard)"
body = slide.shapes.placeholders[1].text_frame
body.text = "Définition : Un outil visuel centralisant tous les KPIs en temps réel pour piloter la campagne et prendre des décisions rapides."
p = body.add_paragraph()
p.text = "Pourquoi c'est important :"
p.level = 1
p = body.add_paragraph()
p.text = "Permet d'ajuster le tir en direct si on voit qu'une publicité ne fonctionne pas."
p.level = 2
p = body.add_paragraph()
p.text = "Exemple :"
p.level = 1
p = body.add_paragraph()
p.text = "Un écran Looker Studio ou un fichier Excel qui montre chaque jour : Dépenses du jour, Ventes du jour, et ROI actuel."
p.level = 2


# ==========================================
# CONCLUSION
# ==========================================
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Conclusion Générale"
body = slide.shapes.placeholders[1].text_frame
body.text = "En synthèse, gérer une campagne de A à Z implique de maîtriser :"
p = body.add_paragraph()
p.text = "1. Les Ressources (C5) : Structurer un budget cohérent et un planning rigoureux."
p.level = 1
p = body.add_paragraph()
p.text = "2. L'Humain (C6) : Savoir s'entourer des bons métiers et répartir les responsabilités."
p.level = 1
p = body.add_paragraph()
p.text = "3. Le Résultat (C7) : Avoir un plan d'action clair, mesurable en direct via des KPIs et un tableau de bord."
p.level = 1
p = body.add_paragraph()
p.text = "=> Le secret du marketing moderne n'est pas l'improvisation, mais l'organisation méthodique !"

prs.save("Sequence_3_Campagne_Marketing_Detaillee.pptx")
print("Detailed Presentation generated successfully!")
