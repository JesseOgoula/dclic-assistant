from pptx import Presentation

prs = Presentation()

# Slide 1: Titre & Introduction
slide_layout = prs.slide_layouts[0] # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Séquence 3 : Gestion d'une campagne marketing"
subtitle.text = "Prenez les commandes de A à Z !\n\n- Objectif principal : S'entraîner à utiliser des méthodes de gestion de projet appliquées au marketing.\n- Méthodologie : Travail de groupe autour d'un ordinateur."

# Slide 2: Les 3 Compétences Clés
slide_layout = prs.slide_layouts[1] # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Ce que nous allons maîtriser"
body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "Cette séquence s'articule autour de 3 piliers (décomposés en 9 modules) :"
p = tf.add_paragraph()
p.text = "C5 : Établir un budget et un planning"
p.level = 1
p = tf.add_paragraph()
p.text = "C6 : Travailler en équipe"
p.level = 1
p = tf.add_paragraph()
p.text = "C7 : Préparer un plan d’action et suivre les indicateurs"
p.level = 1

# Slide 3: Pilier 1
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "C5 - Cadrer sa campagne (Budget & Planning)"
body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "Module 5A : Les méthodes de gestion de projets appliquées au marketing numérique."
p = tf.add_paragraph()
p.text = "Module 5B : Construire le planning et le budget d'une campagne."
p = tf.add_paragraph()
p.text = "Module 5C : Étude de cas pratiques (exemples réels de campagnes)."

# Slide 4: Pilier 2
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "C6 - Travailler en équipe"
body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "Module 6A : Les principes fondamentaux d’organisation d’une équipe marketing."
p = tf.add_paragraph()
p.text = "Module 6B : Panorama des différents métiers nécessaires."
p = tf.add_paragraph()
p.text = "Module 6C : Comment bien répartir les tâches entre les experts ?"

# Slide 5: Pilier 3
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "C7 - Déployer et Mesurer le succès"
body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "Module 7A : Structurer son plan d’action et comprendre l'utilité des indicateurs."
p = tf.add_paragraph()
p.text = "Module 7B : Identifier et choisir les bons indicateurs de suivi (KPIs)."
p = tf.add_paragraph()
p.text = "Module 7C : Créer des tableaux de bord pour piloter les performances."

# Slide 6: Conclusion
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Prêts à lancer votre campagne ?"
body_shape = slide.shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "1 Projet"
p = tf.add_paragraph()
p.text = "3 Compétences"
p = tf.add_paragraph()
p.text = "9 Modules"
p = tf.add_paragraph()
p.text = "À vous de jouer !"

prs.save("Sequence_3_Campagne_Marketing.pptx")
print("Presentation generated successfully!")
